#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Build the MSPR3 NTL WMS soutenance deck (v2) — standalone, original layout.
Derived from the assignment (sujet + grille) and the project deliverables,
NOT from the pre-existing generate_pptx.py / existing pptx.

Output: MSPR3-NTL-WMS-Soutenance-v2.pptx (new file, does not overwrite anything).
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# --------------------------------------------------------------------------
# Palette (indigo + teal — deliberately distinct, clean corporate look)
# --------------------------------------------------------------------------
INK    = RGBColor(0x1E, 0x29, 0x52)   # deep indigo — titles / dark bg
INK2   = RGBColor(0x2A, 0x37, 0x63)   # lighter indigo — inset panels
TEAL   = RGBColor(0x0D, 0x94, 0x88)   # accent
TEAL_L = RGBColor(0x6E, 0xD6, 0xCC)   # light teal — on dark bg
SLATE  = RGBColor(0x33, 0x41, 0x55)   # strong body text
GRAY   = RGBColor(0x64, 0x74, 0x8B)   # secondary text
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)   # card / band fill
LINE   = RGBColor(0xCB, 0xD5, 0xE1)   # borders
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0x9F, 0xB3, 0xC8)   # subtitle on dark
PAPER  = RGBColor(0xD5, 0xDE, 0xEA)   # light text on dark
CODEBG = RGBColor(0x0F, 0x17, 0x2A)

SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.6)
BAR_W  = Inches(0.16)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# --------------------------------------------------------------------------
# Primitives
# --------------------------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, l, t, w, h, fill=None, line=None, line_w=Pt(1.0),
         shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, l, t, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w
    sp.shadow.inherit = False
    return sp


def tb(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def para(tf, runs, first=False, align=PP_ALIGN.LEFT, space_after=6,
         space_before=0, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if line:
        p.line_spacing = line
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, fmt in runs:
        r = p.add_run()
        r.text = text
        f = r.font
        f.size = Pt(fmt.get("size", 14))
        f.bold = fmt.get("bold", False)
        f.italic = fmt.get("italic", False)
        f.name = fmt.get("name", "Calibri")
        f.color.rgb = fmt.get("color", SLATE)
    return p


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# --------------------------------------------------------------------------
# Chrome (content slides): left accent bar, title, underline, footer, number
# --------------------------------------------------------------------------
def chrome(s, title, idx, subtitle=None):
    rect(s, 0, 0, BAR_W, SH, fill=TEAL)
    tf = tb(s, MARGIN, Inches(0.42), SW - 2 * MARGIN, Inches(0.7))
    para(tf, [(title, {"size": 26, "bold": True, "color": INK})], first=True)
    rect(s, MARGIN, Inches(1.16), Inches(1.05), Pt(3), fill=TEAL)
    if subtitle:
        st = tb(s, MARGIN, Inches(1.22), SW - 2 * MARGIN, Inches(0.35))
        para(st, [(subtitle, {"size": 12.5, "italic": True, "color": GRAY})],
             first=True)
    ft = tb(s, MARGIN, SH - Inches(0.42), Inches(7), Inches(0.3))
    para(ft, [("MSPR3 · NordTransit Logistics — WMS-DB",
               {"size": 9, "color": GRAY})], first=True)
    fn = tb(s, SW - Inches(1.3), SH - Inches(0.42), Inches(0.7), Inches(0.3))
    para(fn, [(str(idx), {"size": 10, "bold": True, "color": TEAL})],
         first=True, align=PP_ALIGN.RIGHT)


def bullets_slide(idx, title, items, subtitle=None, note=""):
    """items: list of (level:int, lead:str, rest:str)."""
    s = slide()
    chrome(s, title, idx, subtitle=subtitle)
    top = Inches(1.62) if subtitle else Inches(1.5)
    tf = tb(s, MARGIN + Inches(0.08), top,
            SW - 2 * MARGIN - Inches(0.16), SH - top - Inches(0.6))
    for i, (lvl, lead, rest) in enumerate(items):
        sz = 15 if lvl == 0 else 13
        bullet = "▪  " if lvl == 0 else "–  "
        indent = "" if lvl == 0 else "        "
        runs = [(indent + bullet,
                 {"size": sz, "bold": True, "color": TEAL if lvl == 0 else GRAY})]
        if lead:
            runs.append((lead, {"size": sz, "bold": True,
                                "color": INK if lvl == 0 else SLATE}))
            if rest:
                runs.append(("   " + rest, {"size": sz, "color": SLATE}))
        else:
            runs.append((rest, {"size": sz, "color": SLATE}))
        para(tf, runs, first=(i == 0),
             space_after=(9 if lvl == 0 else 4), line=1.05)
    if note:
        notes(s, note)
    return s


def code_slide(idx, title, lines, note=""):
    s = slide()
    chrome(s, title, idx)
    rect(s, MARGIN, Inches(1.5), SW - 2 * MARGIN, SH - Inches(2.25),
         fill=CODEBG, line=LINE)
    tf = tb(s, MARGIN + Inches(0.28), Inches(1.72),
            SW - 2 * MARGIN - Inches(0.56), SH - Inches(2.6))
    for i, ln in enumerate(lines):
        color = TEAL_L if ln.strip().startswith("--") else PAPER
        para(tf, [(ln if ln else " ",
                   {"size": 12, "name": "Consolas", "color": color})],
             first=(i == 0), space_after=1, line=1.0)
    if note:
        notes(s, note)
    return s


def add_table(s, l, t, w, h, data, col_ratios=None, font=11):
    rows, cols = len(data), len(data[0])
    tbl = s.shapes.add_table(rows, cols, l, t, w, h).table
    if col_ratios:
        total = float(sum(col_ratios))
        for i, r in enumerate(col_ratios):
            tbl.columns[i].width = int(w * r / total)
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            f = r.font
            f.size = Pt(font)
            f.name = "Calibri"
            if ri == 0:
                f.bold = True
                f.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = INK
            else:
                f.color.rgb = SLATE
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
    return tbl


def node(s, l, t, w, h, title, sub, fill=LIGHT, tcolor=INK, border=TEAL):
    rect(s, l, t, w, h, fill=fill, line=border, line_w=Pt(1.5),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = tb(s, l + Inches(0.08), t, w - Inches(0.16), h, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(title, {"size": 13.5, "bold": True, "color": tcolor})],
         first=True, align=PP_ALIGN.CENTER, space_after=1)
    if sub:
        para(tf, [(sub, {"size": 10, "color": GRAY})], align=PP_ALIGN.CENTER)


def conn(s, x1, y1, x2, y2, color=SLATE, w=Pt(1.75)):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = w
    c.shadow.inherit = False
    return c


# ==========================================================================
# S1 — TITLE
# ==========================================================================
s = slide()
rect(s, 0, 0, SW, SH, fill=INK)
rect(s, 0, SH - Inches(0.22), SW, Inches(0.22), fill=TEAL)
rect(s, Inches(0.9), Inches(2.05), Inches(1.4), Pt(4), fill=TEAL)
tf = tb(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(2.1))
para(tf, [("NordTransit Logistics", {"size": 40, "bold": True, "color": WHITE})],
     first=True, space_after=2)
para(tf, [("Industrialisation de la base de données WMS",
           {"size": 25, "color": MUTED})])
tf2 = tb(s, Inches(0.9), Inches(4.55), Inches(11.5), Inches(1.0))
para(tf2, [("MSPR3 — EPSI Nantes · Blocs E6.3 / E6.4 · 2025-2026",
            {"size": 14, "color": WHITE})], first=True, space_after=4)
para(tf2, [("Ianis Puichaud · Blaise Carel · Ojvind · Zaid Abouyaala",
            {"size": 14, "bold": True, "color": TEAL_L})])
rect(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.95), fill=INK2)
etf = tb(s, Inches(1.12), Inches(5.97), Inches(11.05), Inches(0.72),
         anchor=MSO_ANCHOR.MIDDLE)
para(etf, [("EN — ", {"size": 11, "bold": True, "color": TEAL_L}),
           ("Design & industrialization of NTL's WMS relational database: "
            "data model, high availability, disaster recovery "
            "(RTO 1h / RPO 15min), access security and monitoring.",
            {"size": 11, "italic": True, "color": PAPER})], first=True)
notes(s, "Orateur : Ianis (~0:30). Accroche : le WMS est le coeur de NTL, "
         "on a industrialise sa base. Presenter l'equipe et le plan.")

# ==========================================================================
# S2 — CONTEXTE & ENJEU
# ==========================================================================
bullets_slide(2, "Contexte & enjeu", [
    (0, "PME logistique Hauts-de-France",
     "siège Lille + 3 entrepôts (Lens, Valenciennes, Arras) + cross-dock saisonnier"),
    (0, "WMS = cœur métier",
     "sa panne bloque réception & expédition sur les 4 sites (5h30–18h30)"),
    (0, "Contraintes fortes",
     "maintenance nocturne uniquement · équipe IT réduite · liens 200 Mb/s"),
    (0, "Existant recensé",
     "WMS-DB MySQL / Ubuntu 20.04 · NAS RAID5 · sauvegardes sans test ni RTO/RPO"),
    (0, "Enjeu",
     "aujourd'hui, une panne = 4 sites à l'arrêt, sans garantie de reprise"),
], note="Orateur : Ianis (~1:30). Poser vite le vocabulaire (WMS, cross-dock). "
        "Marteler la criticite : 4 sites a l'arret. C'est le 'pourquoi' du projet.")

# ==========================================================================
# S3 — DEMARCHE & PILOTAGE PROJET
# ==========================================================================
bullets_slide(3, "Démarche & pilotage projet", [
    (0, "Équipe 4 pers · 19 h",
     "Ianis + Blaise (données) · Ojvind + Zaid (infra & exploitation)"),
    (0, "Méthode",
     "cadrage → conception → industrialisation → documentation"),
    (0, "Décisions tracées", "6 ADR (le sujet en exige ≥ 3)"),
    (0, "Pilotage", "planning + jalons + suivi de tâches"),
    (0, "Maîtrise des risques",
     "registre de 9 risques (impacts + mesures de mitigation)"),
    (0, "Difficulté (à suivre)",
     "sujet ambigu → on a tranché par un postulat assumé — détail en fin"),
], note="Orateur : Ianis (~1:30). Montrer une demarche maitrisee, pas du "
        "tatonnement. Teaser la difficulte sans la developper ici.")

# ==========================================================================
# S4 — MODELE DE DONNEES (MCD image + points cles)
# ==========================================================================
s = slide()
chrome(s, "Modèle Conceptuel de Données (MCD)", 4,
       subtitle="Point 1 — 7 entités · 8 associations → 8 tables au MLD")
_mcd = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                    "01-architecture-technique", "mcd", "wms-mcd.png")
img_l, img_t, img_w = Inches(0.5), Inches(1.62), Inches(7.55)
rect(s, img_l - Inches(0.06), img_t - Inches(0.06),
     img_w + Inches(0.12), Inches(5.30), fill=WHITE, line=LINE, line_w=Pt(1.0))
s.shapes.add_picture(_mcd, img_l, img_t, width=img_w)
pan_l = img_l + img_w + Inches(0.25)
pan_w = SW - pan_l - MARGIN
rect(s, pan_l, Inches(1.62), pan_w, Inches(5.18), fill=LIGHT, line=LINE,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
ptf = tb(s, pan_l + Inches(0.22), Inches(1.84), pan_w - Inches(0.44), Inches(4.8))
para(ptf, [("Points clés", {"size": 14, "bold": True, "color": TEAL})],
     first=True, space_after=11)
for _lead, _rest in [
    ("7 entités · 8 associations", "→ 8 tables au MLD"),
    ("STOCK", "= article × emplacement × quantité (UNIQUE)"),
    ("mouvement.id_client nullable", "transfert interne = NULL"),
    ("Intégrité par le schéma", "FK · CHECK · UNIQUE · RESTRICT"),
    ("3NF · 0 trigger", "clés INT UNSIGNED"),
]:
    para(ptf, [("▪  ", {"size": 12.5, "bold": True, "color": TEAL}),
               (_lead + " ", {"size": 12.5, "bold": True, "color": INK}),
               (_rest, {"size": 12.5, "color": SLATE})],
         space_after=9, line=1.05)
notes(s, "Orateurs : Blaise + Ianis (~2:00). Montrer le diagramme MCD. "
         "STOCK = article x emplacement x quantite. id_client nullable "
         "(interne = NULL). 4 types de mouvement (CHECK), quantite > 0. "
         "Punch : integrite garantie par le schema, pas par du code fragile.")

# ==========================================================================
# S5 — OPTIMISATION & PERFORMANCE
# ==========================================================================
bullets_slide(5, "Optimisation & performance", [
    (0, "Démarche orientée usages",
     "requêtes les plus fréquentes identifiées → index ciblés"),
    (0, "Requêtes types",
     "stock par site / article · mouvements par période · recherche SKU"),
    (0, "Index",
     "PK / FK + index composites (id_site, id_article) et (date_mouvement)"),
    (0, "Choix physiques",
     "INT UNSIGNED (FK légères) · InnoDB (ACID + FK) · utf8mb4"),
], note="Orateur : Blaise (~1:00). Relier chaque index a un usage reel. "
        "Montrer qu'on a pense les acces avant d'indexer.")

# ==========================================================================
# S6 — CHOIX SGBD
# ==========================================================================
bullets_slide(6, "Choix du SGBD", [
    (0, "MariaDB 11.4 LTS", "support jusqu'en 2029"),
    (0, "Argument n° 1 : drop-in de MySQL",
     "migration du WMS quasi nulle, équipe IT non reformée"),
    (0, "vs PostgreSQL",
     "réécriture du WMS-APP (driver, dialecte, types) + équipe sans "
     "compétence PG → écarté"),
    (0, "Positionnement",
     "« meilleur compromis dans CE contexte », pas « MariaDB est mieux »"),
], note="Orateur : Ojvind (~1:00). Ne JAMAIS dire 'MariaDB est superieur'. "
        "L'argument est contextuel : continuite avec l'existant MySQL.")

# ==========================================================================
# S7 — ARCHITECTURE HA (diagram)
# ==========================================================================
s = slide()
chrome(s, "Architecture Haute Disponibilité", 7,
       subtitle="WMS-APP → HAProxy → Master, réplication vers 2 réplicas")
node(s, Inches(0.85), Inches(3.15), Inches(2.05), Inches(1.05),
     "WMS-APP", "Application")
node(s, Inches(3.30), Inches(3.15), Inches(1.95), Inches(1.05),
     "HAProxy", "Répartiteur")
node(s, Inches(5.60), Inches(3.15), Inches(2.35), Inches(1.05),
     "BDD1 — Master", "écritures", fill=INK, tcolor=WHITE, border=INK)
node(s, Inches(9.15), Inches(1.75), Inches(3.20), Inches(1.05),
     "BDD2 — Slave HA", "réplica local")
node(s, Inches(9.15), Inches(4.55), Inches(3.20), Inches(1.05),
     "BDD3 — Slave PRA", "réplica distant")
conn(s, Inches(2.90), Inches(3.675), Inches(3.30), Inches(3.675))
conn(s, Inches(5.25), Inches(3.675), Inches(5.60), Inches(3.675))
conn(s, Inches(7.95), Inches(3.45), Inches(9.15), Inches(2.275), color=TEAL)
conn(s, Inches(7.95), Inches(3.90), Inches(9.15), Inches(5.075), color=TEAL)
lab = tb(s, Inches(7.55), Inches(2.35), Inches(2.2), Inches(0.3))
para(lab, [("réplication async", {"size": 9, "italic": True, "color": TEAL})],
     first=True)
lab2 = tb(s, Inches(7.55), Inches(4.55), Inches(2.2), Inches(0.3))
para(lab2, [("réplication async", {"size": 9, "italic": True, "color": TEAL})],
     first=True)
cap = tb(s, MARGIN, Inches(6.05), SW - 2 * MARGIN, Inches(0.9))
para(cap, [("Réplication asynchrone + binlogs", {"size": 12.5, "bold": True,
            "color": INK}),
           ("  — faible impact sur le master, adaptée à une équipe IT "
            "réduite.", {"size": 12.5, "color": SLATE})], first=True,
     space_after=2)
para(cap, [("BDD2 couvre une panne service / VM locale · BDD3 couvre la perte "
            "du site de Lille.", {"size": 12.5, "color": SLATE})])
notes(s, "Orateur : Ojvind (~1:30). Deux niveaux de secours : local (HA) et "
         "distant (PRA). Async choisi pour la simplicite ; Galera synchrone = V2.")

# ==========================================================================
# S8 — PRA, SAUVEGARDES & RESTAURATION
# ==========================================================================
bullets_slide(8, "PRA, sauvegardes & restauration", [
    (0, "RTO tenus (cible 1 h)",
     "panne service 5–10 min · perte VM 5–10 min · perte site 20–40 min"),
    (0, "RPO 15 min", "garanti par les binlogs (PITR)"),
    (0, "Sauvegardes automatisées",
     "complète quotidienne (mariabackup) · transactionnel (binlogs) · "
     "rotation / rétention · vérification · notification"),
    (0, "Tests planifiés",
     "bascule HA trimestrielle · PRA semestrielle · restauration mensuelle"),
    (0, "Principe", "une sauvegarde non testée n'est pas une sauvegarde"),
], subtitle="Points 2 & 3 — scripts détaillés en annexe",
   note="Orateur : Zaid (~2:00). Insister : RTO/RPO tenus AVEC marge. "
        "Les scripts existent (annexe). Restauration testee = differenciateur.")

# ==========================================================================
# S9 — SECURITE DES ACCES
# ==========================================================================
bullets_slide(9, "Sécurité des accès", [
    (0, "Moindre privilège",
     "compte app_wms (CRUD tables métier) ≠ compte dba (ALL + admin)"),
    (0, "Pas de root applicatif",
     "l'application n'utilise jamais un compte privilégié"),
    (0, "Renforcement",
     "MFA sur les accès admin · secrets en coffre (jamais en clair)"),
    (0, "Traçabilité", "comptes nominatifs pour l'audit"),
], subtitle="Point 4 du cahier des charges",
   note="Orateur : Zaid (~1:00). Le principe de moindre privilege est explicite "
        "dans le sujet ET la grille. Montrer la separation app / dba.")

# ==========================================================================
# S10 — NOTE CODIR
# ==========================================================================
bullets_slide(10, "Note Comité de direction — risques cyber", [
    (0, "5 risques majeurs",
     "arrêt prolongé · vol d'accès · ransomware + backups KO · "
     "exfiltration données clients · perte de traçabilité"),
    (0, "Langage direction", "impact métier, pas de jargon technique"),
    (0, "Plan priorisé",
     "P1 immédiat (6 sem.) · P2 (3 mois) · P3 (6 mois)"),
    (0, "Conformité", "RGPD : notification d'une fuite sous 72 h"),
], subtitle="Point 3 — favoriser l'adhésion des décideurs",
   note="Orateur : Zaid (~1:30). Registre non technique = competence 'adhesion "
        "decideurs cybercriminalite'. Parler impact business, priorisation.")

# ==========================================================================
# S11 — SUPERVISION
# ==========================================================================
bullets_slide(11, "Supervision", [
    (0, "Stack",
     "Prometheus + Grafana · mysqld_exporter + node_exporter · slow query log"),
    (0, "5 indicateurs + seuils",
     "latence SQL · connexions actives · retard de réplication · disque "
     "libre · succès sauvegardes"),
    (0, "Seuils reliés à un objectif",
     "réplication > 30 s → RPO menacé · disque < 10 % → arrêt "
     "MariaDB possible"),
    (0, "Procédure d'analyse", "définie pour chaque alerte"),
], subtitle="Point 5 — tableau de bord à 5 indicateurs critiques",
   note="Orateur : Ojvind (~1:30). Ne pas reciter les seuils : relier chacun a "
        "un objectif metier (RPO, disponibilite). C'est ce qui distingue.")

# ==========================================================================
# S12 — EXPLOITATION RUNBOOK & LOGS
# ==========================================================================
bullets_slide(12, "Exploitation — RunBook & analyse de logs", [
    (0, "RunBook",
     "start / stop + health checks · checklists jour / semaine · procédure "
     "incident (détection → diagnostic → correction → normal)"),
    (0, "Escalade", "matrice N1 → N4 (15 min → 2 h) + KPIs"),
    (0, "Logs — 5 sources",
     "error.log · slow.log · binlogs · HAProxy · sauvegardes"),
    (0, "Méthode",
     "patterns critiques cartographiés + corrélation entre sources"),
], subtitle="Points 2 & 6 du cahier des charges",
   note="Orateur : Zaid (~1:30). RunBook = exploitabilite au quotidien. "
        "Logs : pas juste lister, montrer patterns + correlation.")

# ==========================================================================
# S13 — DIFFICULTES -> SOLUTIONS
# ==========================================================================
bullets_slide(13, "Difficultés rencontrées → solutions", [
    (0, "Difficulté",
     "sujet ambigu : « PME logistique » + « séparation par client » "
     "= grossiste OU 3PL"),
    (0, "Solution",
     "postulat explicite grossiste (ADR 0003) — itération maîtrisée, "
     "pas tâtonnement"),
    (0, "Conséquences assumées",
     "multi-tenant reporté V2 · 0 trigger · fournisseur en attribut → 8 tables"),
    (0, "Honnêteté",
     "le DDL n'est pas encore exécuté → assumé, c'est notre 1re perspective"),
    (0, "Question attendue du jury",
     "« où est l'isolation par client ? » → attribution (grossiste) "
     "vs isolation (3PL, V2)"),
], subtitle="Le fil rouge de la soutenance",
   note="Orateur : Ianis (~1:30). MOMENT CLE. 'On a tranche une ambiguite par un "
        "postulat assume', JAMAIS 'on n'a pas compris'. Anticiper le piege "
        "isolation/multi-tenant. Rester calme, renvoyer a l'ADR 0003.")

# ==========================================================================
# S14 — RESULTATS & PERSPECTIVES
# ==========================================================================
bullets_slide(14, "Résultats & perspectives V2", [
    (0, "Livré", "8 livrables couvrant tout le cahier des charges"),
    (0, "Chiffres-clés",
     "8 tables 3NF · archi RTO 1h / RPO 15min · 3 serveurs · 5 indicateurs · 6 ADR"),
    (0, "Limites = arbitrages tracés",
     "chaque limite renvoie à un ADR, jamais un oubli"),
    (0, "Perspectives V2",
     "exécution du DDL (prio 1) · multi-tenant si 3PL · lots / FEFO · "
     "expéditions · Galera synchrone"),
], note="Orateur : Blaise (~1:30). Cadrer la V1 comme un MVP, pas un brouillon. "
        "Les perspectives montrent qu'on maitrise le perimetre.")

# ==========================================================================
# S15 — CONCLUSION
# ==========================================================================
s = slide()
rect(s, 0, 0, SW, SH, fill=INK)
rect(s, 0, SH - Inches(0.22), SW, Inches(0.22), fill=TEAL)
rect(s, Inches(0.9), Inches(2.35), Inches(1.4), Pt(4), fill=TEAL)
tf = tb(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.6))
para(tf, [("Une base WMS conçue, sécurisée,",
           {"size": 30, "bold": True, "color": WHITE})], first=True, space_after=2)
para(tf, [("supervisée et résiliente.",
           {"size": 30, "bold": True, "color": WHITE})])
tf2 = tb(s, Inches(0.9), Inches(4.35), Inches(11.5), Inches(0.5))
para(tf2, [("MVP industrialisable — RTO 1h / RPO 15min tenus.",
            {"size": 15, "color": MUTED})], first=True)
tf3 = tb(s, Inches(0.9), Inches(5.15), Inches(11.5), Inches(0.5))
para(tf3, [("« On a tranché une ambiguïté par un postulat assumé. »",
            {"size": 15, "italic": True, "color": TEAL_L})], first=True)
tf4 = tb(s, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.5))
para(tf4, [("Merci — questions ?", {"size": 18, "bold": True, "color": WHITE})],
     first=True)
notes(s, "Toute l'equipe. Phrase finale = le fil rouge. Enchainer sur "
         "l'entretien : garder les annexes pretes.")

# ==========================================================================
# ANNEXES
# ==========================================================================
code_slide("A1", "Annexe — DDL (extrait : table STOCK)", [
    "-- STOCK : 1 article x 1 emplacement x 1 quantite (entite reifiee)",
    "CREATE TABLE stock (",
    "  id_stock        INT UNSIGNED PRIMARY KEY AUTO_INCREMENT,",
    "  id_article      INT UNSIGNED NOT NULL,",
    "  id_localisation INT UNSIGNED NOT NULL,",
    "  quantite        INT NOT NULL,",
    "  CONSTRAINT uq_stock UNIQUE (id_article, id_localisation),",
    "  CONSTRAINT ck_qte   CHECK  (quantite >= 0),",
    "  CONSTRAINT fk_stock_article  FOREIGN KEY (id_article)",
    "      REFERENCES article(id_article)          ON DELETE RESTRICT,",
    "  CONSTRAINT fk_stock_loc      FOREIGN KEY (id_localisation)",
    "      REFERENCES localisation(id_localisation) ON DELETE RESTRICT",
    ") ENGINE=InnoDB DEFAULT CHARSET=utf8mb4;",
], note="Annexe entretien. Montre l'integrite reelle : UNIQUE, CHECK, FK "
        "RESTRICT, InnoDB, utf8mb4. Le DDL complet couvre les 8 tables.")

bullets_slide("A2", "Annexe — MCD (entités & cardinalités)", [
    (0, "SITE (1,N) — (1,1) LOCALISATION", "un site contient N localisations"),
    (0, "ARTICLE / LOCALISATION → STOCK",
     "stock = réification article × emplacement (UNIQUE)"),
    (0, "STOCK (1,N) — (1,1) MOUVEMENT", "un mouvement porte sur un stock"),
    (0, "CLIENT (0,N) — (0,1) MOUVEMENT",
     "sortie client renseignée · mouvement interne = NULL"),
    (0, "CLIENT (0,N) — (0,N) ARTICLE",
     "→ table associative commande"),
    (0, "UTILISATEUR (1,N) — (1,1) MOUVEMENT", "traçabilité de l'auteur"),
], note="Annexe entretien. Justifie les 4 types de mouvement (CHECK) et le "
        "id_client nullable issu de CONCERNE (0,1 - 0,N).")

bullets_slide("A3", "Annexe — Scripts de sauvegarde", [
    (0, "Complète",
     "mariabackup quotidien 02h00 → NAS + copie externalisée (règle 3-2-1)"),
    (0, "Transactionnel", "archivage des binlogs toutes les ~15 min (RPO)"),
    (0, "Rotation / rétention",
     "complètes conservées · binlogs 48 h · purge automatique"),
    (0, "Vérification", "restauration test mensuelle + contrôle checksum"),
    (0, "Notification", "mail / webhook succès ou échec → supervision"),
    (0, "Orchestration", "cron + script idempotent journalisé"),
], note="Annexe entretien. Repond a l'exigence explicite : complete, "
        "incrementale/transactionnelle, rotation, verification, notification.")

s = slide()
chrome(s, "Annexe — Registre des risques", "A4")
add_table(s, MARGIN, Inches(1.55), SW - 2 * MARGIN, Inches(4.9), [
    ["Risque", "Gravité", "Mitigation"],
    ["Arrêt prolongé de la base", "Critique", "HA + PRA testés, RTO 1h"],
    ["Vol d'un accès applicatif", "Critique",
     "Moindre privilège + MFA + rotation secrets"],
    ["Ransomware + backups inutilisables", "Critique",
     "3-2-1 + copies immuables + tests restauration"],
    ["Exfiltration des données clients", "Élevée",
     "Cloisonnement + audit log + RGPD 72h"],
    ["Perte de traçabilité réglementaire", "Élevée",
     "Audit log inviolable 12 mois"],
    ["Réplication rompue", "Moyenne",
     "Supervision du retard de réplication + alerte"],
], col_ratios=[38, 16, 46], font=12)
notes(s, "Annexe entretien. Le registre complet compte 9 risques ; extrait des "
         "6 principaux ici. Chaque risque a un impact et une mitigation.")

s = slide()
chrome(s, "Annexe — Glossaire FR / EN", "A5")
add_table(s, MARGIN, Inches(1.55), SW - 2 * MARGIN, Inches(5.0), [
    ["Français", "English"],
    ["Gestion d'entrepôt (WMS)", "Warehouse Management System"],
    ["Haute disponibilité", "High Availability (HA)"],
    ["Plan de reprise d'activité", "Disaster Recovery Plan (DRP)"],
    ["Sauvegarde / restauration", "Backup / Restore"],
    ["Objectifs de reprise / de perte", "RTO / RPO"],
    ["Moindre privilège", "Least privilege"],
    ["Réplication", "Replication"],
    ["Journaux (logs)", "Logs"],
], col_ratios=[50, 50], font=12.5)
notes(s, "Annexe. Repond a l'exigence de documentation technique en francais ET "
         "en anglais (competence de la grille).")

# --------------------------------------------------------------------------
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "MSPR3-NTL-WMS-Soutenance-v2.pptx")
prs.save(out)
print("OK - saved: " + out)
print("Slides: " + str(len(prs.slides._sldIdLst)))
