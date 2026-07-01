#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
generate_pptx.py
=================
Genere la presentation de soutenance MSPR3 - Base de donnees WMS NTL.

Usage :
    python generate_pptx.py

Sortie :
    MSPR3-NTL-WMS-Soutenance.pptx (14 slides, format 16:9)

Dependance unique : python-pptx (voir requirements.txt)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

# =====================================================================
# 1. PALETTE DE COULEURS (exacte, cf. brief)
# =====================================================================
NAVY = RGBColor(0x1A, 0x3A, 0x5C)          # sections, headers
GOLD = RGBColor(0xC8, 0xA8, 0x4B)          # PK
STEEL_BLUE = RGBColor(0x2B, 0x6C, 0xB0)    # FK
VIOLET = RGBColor(0x6B, 0x46, 0xC1)        # UNIQUE
FOREST_GREEN = RGBColor(0x27, 0x67, 0x49)  # CHECK / NOT NULL
BORDEAUX = RGBColor(0x8B, 0x1A, 0x1A)      # critique / alerte
BROWN_ORANGE = RGBColor(0x7A, 0x3D, 0x00)  # attention
BG_CREAM = RGBColor(0xFA, 0xFA, 0xF7)      # fond des slides
TEXT_MAIN = RGBColor(0x1A, 0x1A, 0x1A)     # texte principal
TEXT_SECOND = RGBColor(0x6B, 0x6B, 0x6B)   # texte secondaire
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Teintes claires derivees (pour fonds de cartes / lignes de tableau)
NAVY_LIGHT = RGBColor(0xDD, 0xE6, 0xEE)
BLUE_LIGHT = RGBColor(0xDC, 0xE9, 0xF7)
GREEN_LIGHT = RGBColor(0xE1, 0xEF, 0xE3)
GREEN_VLIGHT = RGBColor(0xEA, 0xF5, 0xEB)
BORDEAUX_LIGHT = RGBColor(0xF5, 0xDD, 0xDD)
ORANGE_LIGHT = RGBColor(0xFB, 0xE8, 0xD3)
GOLD_LIGHT = RGBColor(0xF7, 0xEE, 0xD9)
VIOLET_LIGHT = RGBColor(0xEA, 0xE2, 0xF7)
GREY_LIGHT = RGBColor(0xE9, 0xE9, 0xE7)

FONT = "Calibri"

# =====================================================================
# 2. GEOMETRIE GLOBALE (16:9)
# =====================================================================
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
BANNER_H = Inches(0.47)          # ~1.2 cm, bandeau decoratif navy
MARGIN = Inches(0.5)
CONTENT_TOP = Inches(1.35)       # debut de zone de contenu (sous le titre)


# =====================================================================
# 3. FONCTIONS UTILITAIRES
# =====================================================================

def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs, bg=BG_CREAM):
    """Ajoute une slide vierge avec une couleur de fond donnee."""
    layout = prs.slide_layouts[6]  # layout vierge
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def add_banner(slide, color=NAVY):
    """Bandeau decoratif navy en haut de slide (~1.2 cm)."""
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, BANNER_H)
    banner.fill.solid()
    banner.fill.fore_color.rgb = color
    banner.line.fill.background()
    banner.shadow.inherit = False
    return banner


def add_gold_edge(slide, position="bottom", thickness=Inches(0.12)):
    """Fine lisiere doree (utilisee sur les slides titre / conclusion)."""
    top = SLIDE_H - thickness if position == "bottom" else 0
    edge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, SLIDE_W, thickness)
    edge.fill.solid()
    edge.fill.fore_color.rgb = GOLD
    edge.line.fill.background()
    edge.shadow.inherit = False
    return edge


def set_text(tf, lines, size=13, color=TEXT_MAIN, bold=False, italic=False,
             align=PP_ALIGN.LEFT, font=FONT, line_spacing=None, space_after=None):
    """Remplit un text_frame avec une liste de (texte) ou (texte, dict overrides)."""
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, tuple):
            text, overrides = item
        else:
            text, overrides = item, {}
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.font.name = overrides.get("font", font)
        p.font.size = Pt(overrides.get("size", size))
        p.font.bold = overrides.get("bold", bold)
        p.font.italic = overrides.get("italic", italic)
        p.font.color.rgb = overrides.get("color", color)
        p.alignment = overrides.get("align", align)
        ls = overrides.get("line_spacing", line_spacing)
        if ls:
            p.line_spacing = ls
        sa = overrides.get("space_after", space_after)
        if sa is not None:
            p.space_after = Pt(sa)
    return tf


def add_text(slide, left, top, width, height, lines, size=13, color=TEXT_MAIN,
             bold=False, italic=False, align=PP_ALIGN.LEFT, font=FONT,
             anchor=None, line_spacing=None, space_after=None):
    """Cree une zone de texte. `lines` peut etre une string ou une liste."""
    if isinstance(lines, str):
        lines = [lines]
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    if anchor:
        tf.vertical_anchor = anchor
    set_text(tf, lines, size=size, color=color, bold=bold, italic=italic,
              align=align, font=font, line_spacing=line_spacing, space_after=space_after)
    return tb


def add_slide_title(slide, text):
    """Titre de slide : 24pt gras navy, sous le bandeau."""
    return add_text(slide, MARGIN, Inches(0.62), SLIDE_W - 2 * MARGIN, Inches(0.6),
                     text, size=24, color=NAVY, bold=True)


def add_rect(slide, left, top, width, height, fill=None, line_color=None,
             line_width=Pt(1), rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    shape.shadow.inherit = False
    shape.text_frame.clear()
    return shape


def add_card(slide, left, top, width, height, title, body_lines, fill=WHITE,
             title_color=NAVY, body_color=TEXT_MAIN, border_color=None,
             title_size=15, body_size=12.5, rounded=True, pad=Inches(0.18)):
    """Carte generique : rectangle rempli + titre + liste de lignes de corps."""
    add_rect(slide, left, top, width, height, fill=fill, line_color=border_color,
              rounded=rounded)
    lines = [(title, {"size": title_size, "bold": True, "color": title_color})]
    for b in body_lines:
        lines.append((b, {"size": body_size, "color": body_color, "space_after": 4}))
    add_text(slide, left + pad, top + pad, width - 2 * pad, height - 2 * pad, lines)


def add_connector(slide, x1, y1, x2, y2, color=NAVY, width=Pt(1.75), dashed=False, arrow=True):
    """Trace un connecteur droit avec fleche optionnelle (utilise pour les schemas)."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = width
    if dashed:
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if arrow:
        ln = conn.line._get_or_add_ln()
        tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
        ln.append(tail)
    return conn


def add_legend_swatch(slide, left, top, color, label, size=11):
    """Petit carre colore + libelle (utilise pour les legendes de code couleur)."""
    sq = add_rect(slide, left, top, Inches(0.16), Inches(0.16), fill=color)
    add_text(slide, left + Inches(0.22), top - Inches(0.04), Inches(1.9), Inches(0.3),
              label, size=size, color=TEXT_MAIN)


def add_mini_table(slide, left, top, width, row_height, header_text, rows, header_fill=NAVY):
    """
    Mini-tableau 2 colonnes pour le MLD (slide 6).
    rows : liste de tuples (texte_ligne, couleur_texte, gras)
    """
    # en-tete
    add_rect(slide, left, top, width, row_height, fill=header_fill, rounded=False)
    add_text(slide, left + Inches(0.06), top, width - Inches(0.12), row_height,
              header_text, size=9.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    y = top + row_height
    for text, color, bold in rows:
        add_rect(slide, left, y, width, row_height, fill=WHITE, line_color=GREY_LIGHT, line_width=Pt(0.5))
        add_text(slide, left + Inches(0.06), y, width - Inches(0.12), row_height,
                  text, size=7.6, bold=bold, color=color, anchor=MSO_ANCHOR.MIDDLE)
        y += row_height
    return y  # bas du tableau


# =====================================================================
# SLIDE 1 - TITRE
# =====================================================================
def build_slide_01(prs):
    slide = blank_slide(prs, bg=NAVY)

    add_text(slide, Inches(1), Inches(2.3), Inches(11.33), Inches(1.4),
              "Base de données WMS — NTL", size=36, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(3.5), Inches(11.33), Inches(0.6),
              "Conception · Haute disponibilité · Industrialisation",
              size=18, color=GOLD, align=PP_ALIGN.CENTER, italic=True)

    add_text(slide, Inches(1), Inches(4.6), Inches(11.33), Inches(0.5),
              "Ianis · Blaise · Ojvind · Zaid",
              size=16, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(5.2), Inches(11.33), Inches(0.5),
              "MSPR3 — EPSI Nantes — Juillet 2026",
              size=12, color=RGBColor(0xC9, 0xD4, 0xDF), align=PP_ALIGN.CENTER)

    add_gold_edge(slide, "bottom")
    return slide


# =====================================================================
# SLIDE 2 - SOMMAIRE
# =====================================================================
def build_slide_02(prs):
    slide = blank_slide(prs)
    add_banner(slide)
    add_slide_title(slide, "Sommaire")

    items = [
        ("1", "Contexte & démarche", STEEL_BLUE),
        ("2", "Modèle de données (MCD / MLD)", GOLD),
        ("3", "Infrastructure technique", NAVY),
        ("4", "Haute Disponibilité & PRA", BROWN_ORANGE),
        ("5", "Sécurité & Supervision", BORDEAUX),
        ("6", "Résultats & Perspectives V2", FOREST_GREEN),
    ]

    cols, rows = 3, 2
    gap = Inches(0.3)
    card_w = (SLIDE_W - 2 * MARGIN - (cols - 1) * gap) / cols
    card_h = Inches(2.15)
    top0 = Inches(2.1)

    for idx, (num, label, color) in enumerate(items):
        c = idx % cols
        r = idx // cols
        left = MARGIN + c * (card_w + gap)
        top = top0 + r * (card_h + gap)
        add_rect(slide, left, top, card_w, card_h, fill=color, rounded=True)
        add_text(slide, left + Inches(0.2), top + Inches(0.15), card_w - Inches(0.4), Inches(0.9),
                  num, size=40, bold=True, color=WHITE)
        add_text(slide, left + Inches(0.2), top + card_h - Inches(0.85), card_w - Inches(0.4), Inches(0.7),
                  label, size=15, bold=True, color=WHITE)
    return slide


# =====================================================================
# SLIDE 3 - CONTEXTE NTL
# =====================================================================
def build_slide_03(prs):
    slide = blank_slide(prs)
    add_banner(slide)
    add_slide_title(slide, "Contexte & Périmètre")

    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2
    top = Inches(1.6)
    height = Inches(3.4)

    add_card(slide, MARGIN, top, col_w, height, "Qui est NTL ?",
              ["PME logistique Hauts-de-France",
               "1 siège social — Lille",
               "3 entrepôts régionaux",
               "WMS propriétaire en place"],
              border_color=NAVY, title_color=NAVY)

    add_card(slide, MARGIN + col_w + Inches(0.3), top, col_w, height, "Notre mission",
              ["Concevoir la BDD du WMS",
               "Haute disponibilité (HA) + PRA",
               "Sécurité des accès & des données",
               "Supervision continue"],
              border_color=STEEL_BLUE, title_color=STEEL_BLUE)

    # 3 chiffres-cles en bas
    stats = [("76h", "équipe"), ("4", "livrables"), ("6", "ADR")]
    kw = (SLIDE_W - 2 * MARGIN - 2 * Inches(0.3)) / 3
    ktop = Inches(5.35)
    for i, (num, label) in enumerate(stats):
        left = MARGIN + i * (kw + Inches(0.3))
        add_rect(slide, left, ktop, kw, Inches(1.55), fill=WHITE, line_color=GOLD, line_width=Pt(1.5), rounded=True)
        add_text(slide, left, ktop + Inches(0.15), kw, Inches(0.8), num, size=34, bold=True,
                  color=GOLD, align=PP_ALIGN.CENTER)
        add_text(slide, left, ktop + Inches(1.0), kw, Inches(0.4), label, size=13,
                  color=TEXT_SECOND, align=PP_ALIGN.CENTER)
    return slide


# =====================================================================
# SLIDE 4 - POSTULAT DE CADRAGE (FIL ROUGE)
# =====================================================================
def build_slide_04(prs):
    slide = blank_slide(prs)
    add_banner(slide)
    add_slide_title(slide, "Postulat de cadrage — ADR 0003")

    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2
    top = Inches(1.55)
    height = Inches(1.35)

    add_card(slide, MARGIN, top, col_w, height, "Problème",
              ["Le sujet est ambigu : NTL = grossiste ou prestataire 3PL ?"],
              fill=BORDEAUX_LIGHT, title_color=BORDEAUX, body_color=TEXT_MAIN,
              title_size=14, body_size=12.5)

    add_card(slide, MARGIN + col_w + Inches(0.3), top, col_w, height, "Décision",
              ["Postulat retenu : NTL = GROSSISTE (stock propre)"],
              fill=GREEN_VLIGHT, title_color=FOREST_GREEN, body_color=TEXT_MAIN,
              title_size=14, body_size=12.5)

    add_text(slide, MARGIN, Inches(3.15), SLIDE_W - 2 * MARGIN, Inches(0.4),
              "Conséquences", size=15, bold=True, color=NAVY)

    conseq = [
        "0 trigger (pas de gestion multi-tenant)",
        "8 tables (suppression article_stock, ajout commande)",
        "FK simples (pas de FK composites id_article + id_client)",
    ]
    cw = (SLIDE_W - 2 * MARGIN - 2 * Inches(0.25)) / 3
    ctop = Inches(3.65)
    for i, txt in enumerate(conseq):
        left = MARGIN + i * (cw + Inches(0.25))
        add_rect(slide, left, ctop, cw, Inches(1.5), fill=WHITE, line_color=STEEL_BLUE, rounded=True)
        add_text(slide, left + Inches(0.15), ctop + Inches(0.15), cw - Inches(0.3), Inches(1.2),
                  txt, size=12.5, color=TEXT_MAIN)

    # message fil rouge
    add_rect(slide, MARGIN, Inches(5.7), SLIDE_W - 2 * MARGIN, Inches(1.0), fill=NAVY, rounded=True)
    add_text(slide, MARGIN + Inches(0.3), Inches(5.7), SLIDE_W - 2 * MARGIN - Inches(0.6), Inches(1.0),
              "Toute limite V1 est un ARBITRAGE TRACÉ, jamais un oubli",
              size=16, bold=True, italic=True, color=WHITE, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)
    return slide


# =====================================================================
# SLIDE 5 - MCD (schema entite-association simplifie)
# =====================================================================
def build_slide_05(prs):
    slide = blank_slide(prs)
    add_banner(slide)
    add_slide_title(slide, "Modèle Conceptuel de Données — V2")

    def entity(cx, cy, label, w=Inches(1.55), h=Inches(0.55)):
        left, top = cx - w / 2, cy - h / 2
        add_rect(slide, left, top, w, h, fill=WHITE, line_color=NAVY, line_width=Pt(1.5), rounded=False)
        add_text(slide, left, top, w, h, label, size=10.5, bold=True, color=NAVY,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return left, top, w, h

    # positions (centres, en pouces)
    SITE = (Inches(1.3), Inches(1.85))
    UTIL = (Inches(11.9), Inches(1.85))
    LOC = (Inches(1.3), Inches(3.15))
    CLIENT = (Inches(11.9), Inches(3.15))
    ARTICLE = (Inches(1.3), Inches(4.55))
    STOCK = (Inches(4.6), Inches(4.55))
    MOUV = (Inches(7.9), Inches(6.15))
    COMMANDE_W, COMMANDE_H = Inches(2.0), Inches(0.85)
    COMMANDE_C = (Inches(6.6), Inches(1.9))

    entity(*SITE, "SITE")
    entity(*UTIL, "UTILISATEUR")
    entity(*LOC, "LOCALISATION")
    entity(*CLIENT, "CLIENT")
    entity(*ARTICLE, "ARTICLE")
    entity(*STOCK, "STOCK")
    entity(*MOUV, "MOUVEMENT", w=Inches(1.7))

    # association N-N COMMANDE (fond or)
    cleft, ctop = COMMANDE_C[0] - COMMANDE_W / 2, COMMANDE_C[1] - COMMANDE_H / 2
    add_rect(slide, cleft, ctop, COMMANDE_W, COMMANDE_H, fill=GOLD, rounded=False)
    add_text(slide, cleft, ctop + Inches(0.05), COMMANDE_W, Inches(0.35),
              "COMMANDE", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, cleft, ctop + Inches(0.38), COMMANDE_W, Inches(0.4),
              "quantite_commandee, date_commande", size=8, color=WHITE, align=PP_ALIGN.CENTER)

    def label_at(x, y, txt, color=TEXT_SECOND):
        add_rect(slide, x - Inches(0.32), y - Inches(0.12), Inches(0.64), Inches(0.24), fill=BG_CREAM)
        add_text(slide, x - Inches(0.32), y - Inches(0.14), Inches(0.64), Inches(0.26),
                  txt, size=8, color=color, align=PP_ALIGN.CENTER, bold=True)

    # liens 1-N (bleu acier)
    add_connector(slide, SITE[0], SITE[1] + Inches(0.28), LOC[0], LOC[1] - Inches(0.28), color=STEEL_BLUE)
    label_at(SITE[0] - Inches(0.35), (SITE[1] + LOC[1]) / 2, "(1,N)-(1,1)")

    add_connector(slide, LOC[0] + Inches(0.6), LOC[1] + Inches(0.15), STOCK[0] - Inches(0.7), STOCK[1] - Inches(0.15), color=STEEL_BLUE)
    label_at((LOC[0] + STOCK[0]) / 2, (LOC[1] + STOCK[1]) / 2 - Inches(0.35), "(1,N)-(1,1)")

    add_connector(slide, ARTICLE[0] + Inches(0.78), ARTICLE[1], STOCK[0] - Inches(0.9), STOCK[1], color=STEEL_BLUE)
    label_at((ARTICLE[0] + STOCK[0]) / 2, ARTICLE[1] - Inches(0.32), "(1,N)-(1,1)")

    add_connector(slide, STOCK[0] + Inches(0.3), STOCK[1] + Inches(0.28), MOUV[0] - Inches(1.0), MOUV[1] - Inches(0.25), color=STEEL_BLUE)
    label_at((STOCK[0] + MOUV[0]) / 2 - Inches(0.6), (STOCK[1] + MOUV[1]) / 2, "(0,N)-(1,1)")

    add_connector(slide, UTIL[0] - Inches(0.5), UTIL[1] + Inches(0.4), MOUV[0] + Inches(0.4), MOUV[1] - Inches(0.3), color=STEEL_BLUE)
    label_at((UTIL[0] + MOUV[0]) / 2 + Inches(0.6), (UTIL[1] + MOUV[1]) / 2 - Inches(0.5), "(0,N)-(1,1)")

    add_connector(slide, UTIL[0], UTIL[1] + Inches(0.28), CLIENT[0], CLIENT[1] - Inches(0.28), color=STEEL_BLUE)
    label_at(UTIL[0] + Inches(0.4), (UTIL[1] + CLIENT[1]) / 2, "(1,N)-(1,1)")

    # lien nullable CLIENT-MOUVEMENT (gris, pointille)
    add_connector(slide, CLIENT[0] - Inches(0.5), CLIENT[1] + Inches(0.35), MOUV[0] + Inches(0.85), MOUV[1] - Inches(0.2),
                   color=TEXT_SECOND, dashed=True)
    label_at((CLIENT[0] + MOUV[0]) / 2 + Inches(0.9), (CLIENT[1] + MOUV[1]) / 2 - Inches(0.1), "(0,N)-(0,1)", color=TEXT_SECOND)

    # liens N-N vers COMMANDE (or)
    add_connector(slide, CLIENT[0] - Inches(0.78), CLIENT[1] - Inches(0.4), COMMANDE_C[0] + COMMANDE_W / 2, COMMANDE_C[1] + Inches(0.1), color=GOLD)
    add_connector(slide, ARTICLE[0] + Inches(0.5), ARTICLE[1] - Inches(0.4), COMMANDE_C[0] - COMMANDE_W / 2, COMMANDE_C[1] + Inches(0.1), color=GOLD)

    # legende
    add_legend_swatch(slide, MARGIN, Inches(6.85), WHITE, "Entité (bordure navy)")
    add_legend_swatch(slide, MARGIN + Inches(2.6), Inches(6.85), GOLD, "Association N-N")
    add_legend_swatch(slide, MARGIN + Inches(5.0), Inches(6.85), STEEL_BLUE, "Lien FK / cardinalité")
    return slide


# =====================================================================
# SLIDE 6 - MLD PHYSIQUE (8 tables, dense)
# =====================================================================
def build_slide_06(prs):
    slide = blank_slide(prs)
    add_banner(slide)
    add_slide_title(slide, "Modèle Logique & Contraintes — 8 tables")

    # legende couleur
    lx = MARGIN
    ly = Inches(1.18)
    add_legend_swatch(slide, lx, ly, GOLD, "PK", size=10)
    add_legend_swatch(slide, lx + Inches(1.1), ly, STEEL_BLUE, "FK", size=10)
    add_legend_swatch(slide, lx + Inches(2.2), ly, VIOLET, "UNIQUE", size=10)
    add_legend_swatch(slide, lx + Inches(3.5), ly, FOREST_GREEN, "CHECK / NOT NULL", size=10)
    add_legend_swatch(slide, lx + Inches(5.6), ly, TEXT_SECOND, "NULL", size=10)

    tables = {
        "site": [
            ("id_site | INT UNSIGNED · PK", GOLD, True),
            ("libellé | VARCHAR(100) · NOT NULL", TEXT_MAIN, False),
            ("adresse | VARCHAR(255) · NOT NULL", TEXT_MAIN, False),
        ],
        "localisation": [
            ("id_localisation | INT UNSIGNED · PK", GOLD, True),
            ("code | VARCHAR(30) · UNIQUE", VIOLET, True),
            ("zone, allee, etage, place | VARCHAR · NOT NULL", TEXT_MAIN, False),
            ("id_site | INT UNSIGNED · FK→site", STEEL_BLUE, True),
        ],
        "utilisateur": [
            ("id_utilisateur | INT UNSIGNED · PK", GOLD, True),
            ("nom | VARCHAR(100) · NOT NULL", TEXT_MAIN, False),
            ("rôle | VARCHAR(30) · NOT NULL", TEXT_MAIN, False),
        ],
        "client": [
            ("id_client | INT UNSIGNED · PK", GOLD, True),
            ("nom | VARCHAR(100) · NOT NULL", TEXT_MAIN, False),
            ("siret | CHAR(14) · UNIQUE", VIOLET, True),
            ("telephone | VARCHAR(20) · NULL", TEXT_SECOND, False),
            ("status | VARCHAR(20) · DEFAULT 'actif'", TEXT_MAIN, False),
            ("id_utilisateur | INT UNSIGNED · FK→utilisateur", STEEL_BLUE, True),
        ],
        "article": [
            ("id_article | INT UNSIGNED · PK", GOLD, True),
            ("nom | VARCHAR(150) · NOT NULL", TEXT_MAIN, False),
            ("poids | DECIMAL(10,3) · CHECK>0", FOREST_GREEN, True),
            ("fournisseur | VARCHAR(100) · NULL", TEXT_SECOND, False),
            ("type | VARCHAR(50) · NOT NULL", TEXT_MAIN, False),
        ],
        "stock ★": [
            ("id_stock | INT UNSIGNED · PK", GOLD, True),
            ("quantite | INT UNSIGNED · DEFAULT 0", TEXT_MAIN, False),
            ("id_article | INT UNSIGNED · FK→article", STEEL_BLUE, True),
            ("id_localisation | INT UNSIGNED · FK→localisation", STEEL_BLUE, True),
            ("UNIQUE(id_article, id_localisation)", VIOLET, True),
        ],
        "mouvement ★": [
            ("id_mouvement | INT UNSIGNED · PK", GOLD, True),
            ("type | VARCHAR(20) · CHECK IN(4 valeurs)", FOREST_GREEN, True),
            ("reference | VARCHAR(50) · UNIQUE", VIOLET, True),
            ("date | DATE · NOT NULL", TEXT_MAIN, False),
            ("heure | TIME · NOT NULL", TEXT_MAIN, False),
            ("quantite | INT UNSIGNED · CHECK>0", FOREST_GREEN, True),
            ("id_stock | INT UNSIGNED · FK→stock", STEEL_BLUE, True),
            ("id_utilisateur | INT UNSIGNED · FK→utilisateur", STEEL_BLUE, True),
            ("id_client | INT UNSIGNED · NULL · FK→client", TEXT_SECOND, True),
        ],
        "commande": [
            ("id_commande | INT UNSIGNED · PK", GOLD, True),
            ("id_client | INT UNSIGNED · FK→client", STEEL_BLUE, True),
            ("id_article | INT UNSIGNED · FK→article", STEEL_BLUE, True),
            ("quantite_commandee | INT UNSIGNED · CHECK>0", FOREST_GREEN, True),
            ("date_commande | DATETIME · DEFAULT CURRENT_TIMESTAMP", TEXT_MAIN, False),
        ],
    }

    names = list(tables.keys())
    cols = 4
    gap = Inches(0.15)
    card_w = (SLIDE_W - 2 * MARGIN - (cols - 1) * gap) / cols
    row_h = Inches(0.235)
    top0 = Inches(1.55)
    row_gap = Inches(0.15)
    card_block_h = Inches(2.42)

    for idx, name in enumerate(names):
        c = idx % cols
        r = idx // cols
        left = MARGIN + c * (card_w + gap)
        top = top0 + r * (card_block_h + row_gap)
        add_mini_table(slide, left, top, card_w, row_h, name.upper(), tables[name])

    add_text(slide, MARGIN, Inches(7.1), SLIDE_W - 2 * MARGIN, Inches(0.3),
              "Moteur : InnoDB · Charset : utf8mb4 · ON DELETE RESTRICT partout",
              size=10.5, italic=True, color=TEXT_SECOND, align=PP_ALIGN.CENTER)
    return slide


# =====================================================================
# SLIDE 7 - INFRASTRUCTURE TECHNIQUE
# =====================================================================
def build_slide_07(prs):
    slide = blank_slide(prs)
    add_banner(slide)
    add_slide_title(slide, "Infrastructure technique — 3 serveurs MariaDB")

    servers = [
        ("BDD1 — MASTER", NAVY_LIGHT, NAVY, [
            "Rôle : Production, toutes les écritures",
            "OS : Ubuntu Server 22.04 LTS",
            "SGBD : MariaDB 11.4 LTS (jusqu'en 2029)",
            "Config : 2-4 vCPU · 8-16 Go RAM · SSD",
            "Réseau : 1 Gbit/s · site Lille siège",
        ]),
        ("BDD2 — SLAVE HA", BLUE_LIGHT, STEEL_BLUE, [
            "Rôle : Réplication locale, haute dispo",
            "OS : Ubuntu Server 22.04 LTS",
            "SGBD : MariaDB 11.4 LTS",
            "Config : 2-4 vCPU · 8-16 Go RAM · SSD",
            "Réseau : même LAN que BDD1",
        ]),
        ("BDD3 — SLAVE PRA", GREEN_VLIGHT, FOREST_GREEN, [
            "Rôle : Réplication distante, PRA",
            "OS : Ubuntu Server 22.04 LTS",
            "SGBD : MariaDB 11.4 LTS",
            "Config : 2-4 vCPU · 8-16 Go RAM",
            "Réseau : site distant (hors Lille)",
        ]),
    ]

    gap = Inches(0.3)
    col_w = (SLIDE_W - 2 * MARGIN - 2 * gap) / 3
    top = Inches(1.55)
    height = Inches(3.9)

    for i, (title, fill, border, lines) in enumerate(servers):
        left = MARGIN + i * (col_w + gap)
        add_card(slide, left, top, col_w, height, title, lines, fill=fill,
                  title_color=border, border_color=border, title_size=15, body_size=11.5)

    ftop = Inches(5.75)
    add_rect(slide, MARGIN, ftop, SLIDE_W - 2 * MARGIN, Inches(1.35), fill=WHITE, line_color=NAVY, rounded=True)
    add_text(slide, MARGIN + Inches(0.25), ftop + Inches(0.12), SLIDE_W - 2 * MARGIN - Inches(0.5), Inches(0.5),
              "HAProxy 2.x · Prometheus + mysqld_exporter + node_exporter · mariabackup",
              size=12.5, bold=True, color=NAVY)
    add_text(slide, MARGIN + Inches(0.25), ftop + Inches(0.68), SLIDE_W - 2 * MARGIN - Inches(0.5), Inches(0.5),
              "Paramètres clés MariaDB : innodb_buffer_pool_size=60%RAM · sync_binlog=1 · slow_query_log=ON",
              size=11.5, color=TEXT_SECOND)
    return slide


# =====================================================================
# SLIDE 8 - TOPOLOGIE HA
# =====================================================================
def build_slide_08(prs):
    slide = blank_slide(prs)
    add_banner(slide)
    add_slide_title(slide, "Topologie Haute Disponibilité")

    # colonne gauche : flux vertical
    box_w = Inches(3.4)
    cx = MARGIN + box_w / 2

    def flow_box(top, label, fill, border):
        add_rect(slide, cx - box_w / 2, top, box_w, Inches(0.7), fill=fill, line_color=border,
                  line_width=Pt(1.5), rounded=True)
        add_text(slide, cx - box_w / 2, top, box_w, Inches(0.7), label, size=13, bold=True,
                  color=border, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    y1, y2, y3, y4 = Inches(1.6), Inches(2.75), Inches(3.9), Inches(5.35)
    flow_box(y1, "WMS-APP", WHITE, NAVY)
    flow_box(y2, "HAProxy\n(load balancer / health check)", BLUE_LIGHT, STEEL_BLUE)
    flow_box(y3, "BDD1 — MASTER", NAVY_LIGHT, NAVY)

    add_rect(slide, cx - Inches(1.5), y4, Inches(1.4), Inches(0.7), fill=BLUE_LIGHT, line_color=STEEL_BLUE, rounded=True)
    add_text(slide, cx - Inches(1.5), y4, Inches(1.4), Inches(0.7), "BDD2\nSLAVE HA", size=11.5, bold=True,
              color=STEEL_BLUE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, cx + Inches(0.1), y4, Inches(1.4), Inches(0.7), fill=GREEN_VLIGHT, line_color=FOREST_GREEN, rounded=True)
    add_text(slide, cx + Inches(0.1), y4, Inches(1.4), Inches(0.7), "BDD3\nSLAVE PRA", size=11.5, bold=True,
              color=FOREST_GREEN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_connector(slide, cx, y1 + Inches(0.7), cx, y2, color=NAVY, width=Pt(2.25))
    add_connector(slide, cx, y2 + Inches(0.7), cx, y3, color=NAVY, width=Pt(2.25))
    add_connector(slide, cx, y3 + Inches(0.7), cx - Inches(0.8), y4, color=STEEL_BLUE, dashed=True)
    add_connector(slide, cx, y3 + Inches(0.7), cx + Inches(0.8), y4, color=STEEL_BLUE, dashed=True)

    # legende fleches
    ltop = Inches(6.4)
    add_connector(slide, MARGIN, ltop + Inches(0.1), MARGIN + Inches(0.5), ltop + Inches(0.1), color=NAVY, width=Pt(2.25))
    add_text(slide, MARGIN + Inches(0.6), ltop - Inches(0.05), Inches(2.2), Inches(0.3), "Écriture", size=10.5, color=TEXT_MAIN)
    add_connector(slide, MARGIN + Inches(2.6), ltop + Inches(0.1), MARGIN + Inches(3.1), ltop + Inches(0.1), color=STEEL_BLUE, dashed=True)
    add_text(slide, MARGIN + Inches(3.2), ltop - Inches(0.05), Inches(2.5), Inches(0.3), "Réplication (binlogs async)", size=10.5, color=TEXT_MAIN)

    # colonne droite : comptes MariaDB
    rleft = Inches(7.6)
    rw = SLIDE_W - MARGIN - rleft
    add_card(slide, rleft, Inches(1.6), rw, Inches(4.5), "Comptes MariaDB", [
        "appuser : SELECT/INSERT/UPDATE/DELETE, limité par IP",
        "replicator : REPLICATION SLAVE",
        "dba : ALL PRIVILEGES + SSH + MFA",
        "backup : SELECT + LOCK TABLES",
    ], border_color=NAVY, title_size=16, body_size=13)
    return slide


# =====================================================================
# SLIDE 9 - PRA & SCENARIOS DE BASCULE
# =====================================================================
def build_slide_09(prs):
    slide = blank_slide(prs)
    add_banner(slide)
    add_slide_title(slide, "Plan de Reprise d'Activité — RTO 1h / RPO 15min")

    rows = [
        ("Scenario", "Action", "Duree", NAVY, WHITE, True),
        ("Panne service MariaDB BDD1", "Bascule HAProxy → BDD2", "5-10 min", GREEN_VLIGHT, FOREST_GREEN, False),
        ("Perte VM BDD1", "Promotion BDD2 en master", "5-10 min", GREEN_VLIGHT, FOREST_GREEN, False),
        ("Perte site Lille", "Bascule BDD3 PRA", "20-40 min", ORANGE_LIGHT, BROWN_ORANGE, False),
        ("Corruption logique", "Dump + relecture binlogs", "selon volume", BORDEAUX_LIGHT, BORDEAUX, False),
    ]

    top = Inches(1.55)
    row_h = Inches(0.55)
    col_w = [Inches(4.6), Inches(4.9), Inches(2.83)]
    for r, (c1, c2, c3, fill, txt_color, is_header) in enumerate(rows):
        left = MARGIN
        y = top + r * row_h
        for w, txt in zip(col_w, (c1, c2, c3)):
            add_rect(slide, left, y, w, row_h, fill=fill, line_color=WHITE, line_width=Pt(1))
            add_text(slide, left + Inches(0.1), y, w - Inches(0.2), row_h, txt,
                      size=12.5 if not is_header else 13, bold=is_header, color=txt_color,
                      anchor=MSO_ANCHOR.MIDDLE)
            left += w

    ctop = Inches(4.6)
    add_card(slide, MARGIN, ctop, Inches(6.15), Inches(1.35), "Conditions avant bascule HA", [
        "Slave_IO_Running = Yes",
        "Slave_SQL_Running = Yes",
        "Seconds_Behind_Master = 0",
    ], border_color=STEEL_BLUE, title_size=13.5, body_size=12)

    add_card(slide, MARGIN + Inches(6.45), ctop, Inches(6.38), Inches(1.35), "Plan de sauvegardes", [
        "Complète quotidienne 02h00 (mariabackup)",
        "Binlogs conservés 48h · Rétention 14 jours",
        "Externalisation des sauvegardes",
    ], border_color=GOLD, title_size=13.5, body_size=12)
    return slide


# =====================================================================
# SLIDE 10 - SECURITE NOTE CODIR
# =====================================================================
def build_slide_10(prs):
    slide = blank_slide(prs)
    add_banner(slide)
    add_slide_title(slide, "Sécurité — Note de Direction CODIR")

    risks = [
        ("1", "Arrêt BDD (4 sites bloqués)", "CRITIQUE", BORDEAUX_LIGHT, BORDEAUX),
        ("2", "Vol accès applicatif", "CRITIQUE", BORDEAUX_LIGHT, BORDEAUX),
        ("3", "Ransomware + sauvegardes compromises", "CRITIQUE", BORDEAUX_LIGHT, BORDEAUX),
        ("4", "Exfiltration données clients", "ÉLEVÉE", ORANGE_LIGHT, BROWN_ORANGE),
        ("5", "Perte traçabilité réglementaire", "ÉLEVÉE", ORANGE_LIGHT, BROWN_ORANGE),
    ]
    top = Inches(1.5)
    row_h = Inches(0.4)
    for i, (num, label, level, fill, color) in enumerate(risks):
        y = top + i * row_h
        add_rect(slide, MARGIN, y, Inches(9.5), row_h, fill=fill, line_color=WHITE, line_width=Pt(0.75))
        add_text(slide, MARGIN + Inches(0.15), y, Inches(8.2), row_h, f"{num} · {label}",
                  size=12, bold=True, color=TEXT_MAIN, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, MARGIN + Inches(8.3), y, Inches(1.1), row_h, level, size=11.5, bold=True,
                  color=color, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, MARGIN, Inches(3.7), Inches(9.5), Inches(0.3),
              "RGPD : notification de fuite sous 72h", size=11.5, italic=True, color=TEXT_SECOND)

    phases = [
        ("P1 — 6 semaines", ["Sauvegardes 3-2-1", "Cluster HA", "MFA admin", "Coffre a secrets"], BORDEAUX),
        ("P2 — 3 mois", ["Moindre privilège", "Audit log 12 mois"], BROWN_ORANGE),
        ("P3 — 6 mois", ["Exercice de crise", "Assurance cyber"], FOREST_GREEN),
    ]
    gap = Inches(0.25)
    col_w = (SLIDE_W - 2 * MARGIN - 2 * gap) / 3
    ptop = Inches(4.2)
    for i, (title, items, color) in enumerate(phases):
        left = MARGIN + i * (col_w + gap)
        add_card(slide, left, ptop, col_w, Inches(2.5), title, items, border_color=color,
                  title_color=color, title_size=13.5, body_size=12)
    return slide


# =====================================================================
# SLIDE 11 - SUPERVISION
# =====================================================================
def build_slide_11(prs):
    slide = blank_slide(prs)
    add_banner(slide)
    add_slide_title(slide, "Supervision — Prometheus + Grafana")

    add_text(slide, MARGIN, Inches(1.45), SLIDE_W - 2 * MARGIN, Inches(0.4),
              "Stack : Prometheus · Grafana · mysqld_exporter · node_exporter · Slow Query Log",
              size=13, bold=True, color=NAVY)

    headers = ["Indicateur", "Normal", "Alerte", "Critique"]
    data = [
        ("Latence SQL", "<50ms", ">150ms", ">300ms"),
        ("Connexions", "<50", ">100", ">150"),
        ("Réplication (retard)", "0-5s", ">15s", ">30s"),
        ("Espace disque", ">20%", "<15%", "<10%"),
        ("Sauvegardes", "100%", "1 echec", ">1 echec"),
    ]
    col_w = [Inches(3.9), Inches(3.1), Inches(3.1), Inches(2.23)]
    top = Inches(2.05)
    row_h = Inches(0.55)

    left = MARGIN
    for w, h in zip(col_w, headers):
        add_rect(slide, left, top, w, row_h, fill=NAVY)
        add_text(slide, left, top, w, row_h, h, size=12.5, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        left += w

    for r, (label, ok, warn, crit) in enumerate(data):
        y = top + row_h * (r + 1)
        left = MARGIN
        cells = [(label, WHITE, TEXT_MAIN), (ok, GREEN_VLIGHT, FOREST_GREEN),
                 (warn, ORANGE_LIGHT, BROWN_ORANGE), (crit, BORDEAUX_LIGHT, BORDEAUX)]
        for w, (txt, fill, color) in zip(col_w, cells):
            add_rect(slide, left, y, w, row_h, fill=fill, line_color=GREY_LIGHT, line_width=Pt(0.5))
            add_text(slide, left, y, w, row_h, txt, size=12, bold=(fill != WHITE), color=color,
                      align=PP_ALIGN.CENTER if fill != WHITE else PP_ALIGN.LEFT,
                      anchor=MSO_ANCHOR.MIDDLE)
            left += w

    add_text(slide, MARGIN, Inches(5.6), SLIDE_W - 2 * MARGIN, Inches(0.5),
              "Retard de réplication > 30s → RPO 15min menacé",
              size=13, italic=True, bold=True, color=BORDEAUX, align=PP_ALIGN.CENTER)
    return slide


# =====================================================================
# SLIDE 12 - EXPLOITATION RUNBOOK
# =====================================================================
def build_slide_12(prs):
    slide = blank_slide(prs)
    add_banner(slide)
    add_slide_title(slide, "Exploitation — RunBook & Logs")

    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2
    top = Inches(1.5)

    add_card(slide, MARGIN, top, col_w, Inches(5.4), "RunBook", [
        "Procédures : Start/Stop MariaDB, HAProxy, health check",
        "Checklists : quotidienne + hebdomadaire",
        "Matrice d'escalade :",
        "   N1 (15 min) → N2 (30 min) → N3 (1h) → N4 DBA (2h)",
        "KPI d'exploitation",
    ], border_color=NAVY, title_size=16, body_size=13)

    logs = [
        ("error.log", "InnoDB corruption", "critique", BORDEAUX_LIGHT, BORDEAUX),
        ("slow.log", "requêtes lentes", "attention", ORANGE_LIGHT, BROWN_ORANGE),
        ("binlogs", "PITR + replication", "info", BLUE_LIGHT, STEEL_BLUE),
        ("HAProxy", "Backend DOWN", "critique", BORDEAUX_LIGHT, BORDEAUX),
        ("sauvegardes", "Backup failed", "critique", BORDEAUX_LIGHT, BORDEAUX),
    ]
    rleft = MARGIN + col_w + Inches(0.3)
    add_text(slide, rleft, top, col_w, Inches(0.4), "Logs — 5 sources", size=16, bold=True, color=NAVY)
    row_h = Inches(0.85)
    rtop = top + Inches(0.5)
    for i, (name, desc, level, fill, color) in enumerate(logs):
        y = rtop + i * (row_h + Inches(0.08))
        add_rect(slide, rleft, y, col_w, row_h, fill=fill, rounded=True)
        add_text(slide, rleft + Inches(0.2), y + Inches(0.08), col_w - Inches(0.4), Inches(0.35),
                  name, size=13, bold=True, color=color)
        add_text(slide, rleft + Inches(0.2), y + Inches(0.42), col_w - Inches(1.3), Inches(0.35),
                  desc, size=11.5, color=TEXT_MAIN)
        add_text(slide, rleft + col_w - Inches(1.15), y + Inches(0.42), Inches(1.0), Inches(0.35),
                  level, size=10.5, bold=True, color=color, align=PP_ALIGN.RIGHT)
    return slide


# =====================================================================
# SLIDE 13 - RESULTATS & PERSPECTIVES V2
# =====================================================================
def build_slide_13(prs):
    slide = blank_slide(prs)
    add_banner(slide)
    add_slide_title(slide, "Résultats & Perspectives")

    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2
    top = Inches(1.55)
    height = Inches(5.3)

    add_card(slide, MARGIN, top, col_w, height, "Ce qu'on a livré", [
        "✓ MCD + MLD + DDL (8 tables, 3NF)",
        "✓ Architecture HA/PRA (RTO 1h / RPO 15min)",
        "✓ Note de direction sécurité",
        "✓ Stack supervision Prometheus/Grafana",
        "✓ RunBook exploitation complet",
        "✓ 6 ADR traces (sujet exige ≥ 3)",
    ], fill=GREEN_VLIGHT, title_color=FOREST_GREEN, body_size=13.5, title_size=17)

    add_card(slide, MARGIN + col_w + Inches(0.3), top, col_w, height, "V2 — Perspectives", [
        "→ Exécution DDL sur MariaDB 11.4 (1ère priorité)",
        "→ Lots / DLC / FEFO",
        "→ Galera Cluster (HA synchrone)",
        "→ Pivot 3PL si evolution NTL",
        "→ Integration ERP fournisseurs",
    ], fill=BLUE_LIGHT, title_color=STEEL_BLUE, body_size=13.5, title_size=17)
    return slide


# =====================================================================
# SLIDE 14 - CONCLUSION
# =====================================================================
def build_slide_14(prs):
    slide = blank_slide(prs, bg=NAVY)

    lines = [
        "Un postulat assumé, pas une incompréhension.",
        "Un MVP cohérent, documenté, exploitable dès J+1.",
        "6 ADR, 8 tables, RTO 1h — objectifs atteints.",
    ]
    add_text(slide, Inches(1), Inches(1.7), Inches(11.33), Inches(0.4), "Conclusion",
              size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    y = Inches(2.6)
    for line in lines:
        add_text(slide, Inches(1), y, Inches(11.33), Inches(0.7), line, size=20, bold=True,
                  color=WHITE, align=PP_ALIGN.CENTER)
        y += Inches(1.0)

    add_text(slide, Inches(1), Inches(6.2), Inches(11.33), Inches(0.6),
              "Merci · Questions ?", size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    add_gold_edge(slide, "bottom")
    return slide


# =====================================================================
# MAIN
# =====================================================================
def main():
    prs = new_presentation()

    builders = [
        build_slide_01, build_slide_02, build_slide_03, build_slide_04,
        build_slide_05, build_slide_06, build_slide_07, build_slide_08,
        build_slide_09, build_slide_10, build_slide_11, build_slide_12,
        build_slide_13, build_slide_14,
    ]
    for builder in builders:
        builder(prs)

    output_path = "MSPR3-NTL-WMS-Soutenance.pptx"
    prs.save(output_path)
    print(f"Présentation générée : {output_path} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
